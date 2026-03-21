import csv
from pathlib import Path
from typing import List, Dict

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def _chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> List[str]:
    cleaned = " ".join(text.split())
    if not cleaned:
        return []

    chunks: List[str] = []
    start = 0
    while start < len(cleaned):
        end = min(len(cleaned), start + chunk_size)
        chunks.append(cleaned[start:end])
        if end == len(cleaned):
            break
        start = max(0, end - overlap)
    return chunks


def parse_file(path: str, filename: str) -> List[Dict[str, str | int]]:
    suffix = Path(filename).suffix.lower()
    records: List[Dict[str, str | int]] = []

    if suffix == ".pdf":
        reader = PdfReader(path)
        for index, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            for part in _chunk_text(page_text):
                records.append({"text": part, "page": index, "source": filename})
        return records

    if suffix == ".pptx":
        presentation = Presentation(path)
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            for part in _chunk_text("\n".join(texts)):
                records.append({"text": part, "page": index, "source": filename})
        return records

    if suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell) for cell in row if cell is not None]
                if values:
                    rows.append(" | ".join(values))
            for part in _chunk_text("\n".join(rows)):
                records.append({"text": part, "page": 1, "source": f"{filename}:{sheet.title}"})
        return records

    if suffix == ".docx":
        document = DocxDocument(path)
        paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        for part in _chunk_text("\n".join(paragraphs)):
            records.append({"text": part, "page": 1, "source": filename})
        return records

    if suffix == ".csv":
        rows: List[str] = []
        with open(path, "r", encoding="utf-8", newline="") as handle:
            reader = csv.reader(handle)
            for row in reader:
                values = [cell.strip() for cell in row if cell and cell.strip()]
                if values:
                    rows.append(" | ".join(values))
        for part in _chunk_text("\n".join(rows)):
            records.append({"text": part, "page": 1, "source": filename})
        return records

    raise ValueError("Unsupported file type")
